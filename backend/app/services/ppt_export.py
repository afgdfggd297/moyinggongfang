"""HTML 转 PPTX 导出服务 - html2ppt 逐页截图合并版"""
import io
import re
import shutil
import logging
import subprocess
import tempfile
from pathlib import Path
from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.util import Inches

from app.core.config import get_settings

settings = get_settings()
logger = logging.getLogger(__name__)


class PPTExportService:
    """HTML转PPTX - 逐页 html2ppt 截图 + python-pptx 合并"""

    def __init__(self):
        self.upload_dir = Path(settings.UPLOAD_DIR)
        self.upload_dir.mkdir(parents=True, exist_ok=True)
        self._html2ppt_bin = (
            settings.PPT_HTML2PPT_BIN
            or shutil.which("html2ppt")
            or shutil.which("html2ppt.cmd")
            or ""
        )

    def html_to_pptx_sync(self, html_content: str, plan_id: str) -> str:
        """同步导出：每张幻灯片单独截图，再合并到 PPTX"""
        soup = BeautifulSoup(html_content, "lxml")

        all_divs = soup.find_all("div", class_=re.compile(r"\bslide\b"))
        slides = [
            d for d in all_divs
            if "slide" in d.get("class", [])
            and "slide-title" not in d.get("class", [])
            and "slide-num" not in d.get("class", [])
        ]

        if not slides:
            raise ValueError("未找到幻灯片元素")

        logger.info("[PPTX] 找到 %d 个幻灯片，逐页截图...", len(slides))

        screenshots = []

        for i, slide_div in enumerate(slides):
            single_html = self._build_single_slide_html(slide_div, html_content)

            with tempfile.NamedTemporaryFile(
                suffix=".html", delete=False, mode="w", encoding="utf-8"
            ) as f:
                f.write(single_html)
                tmp_html = f.name

            tmp_png = tmp_html.replace(".html", ".png")

            try:
                cmd = [
                    self._html2ppt_bin,
                    tmp_html,
                    "-o", tmp_png.replace(".png", ".pptx"),
                    "--width", str(settings.PPT_PX_W),
                    "--height", str(settings.PPT_PX_H),
                    "--scale", str(settings.PPT_SCALE),
                    "--no-full-page",
                ]
                result = subprocess.run(
                    cmd, capture_output=True, text=True,
                    timeout=settings.PPT_TIMEOUT, shell=True,
                )
                if result.returncode != 0:
                    logger.warning("[PPTX] 第 %d 页 html2ppt 失败: %s", i + 1, result.stderr)
                    continue

                png_path = Path(tmp_png)
                if png_path.exists():
                    screenshots.append(png_path.read_bytes())
                    logger.info("[PPTX] 截图 %d/%d 完成", i + 1, len(slides))
                else:
                    logger.warning("[PPTX] 第 %d 页 PNG 未生成", i + 1)

            finally:
                Path(tmp_html).unlink(missing_ok=True)
                Path(tmp_png).unlink(missing_ok=True)
                Path(tmp_png.replace(".png", ".pptx")).unlink(missing_ok=True)

        if not screenshots:
            raise RuntimeError("所有幻灯片截图失败")

        prs = Presentation()
        prs.slide_width = Inches(settings.PPT_SLIDE_W)
        prs.slide_height = Inches(settings.PPT_SLIDE_H)

        for img_bytes in screenshots:
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            slide.shapes.add_picture(
                io.BytesIO(img_bytes),
                Inches(0), Inches(0),
                Inches(settings.PPT_SLIDE_W), Inches(settings.PPT_SLIDE_H),
            )

        filename = f"ppt_{plan_id}.pptx"
        filepath = self.upload_dir / filename
        prs.save(str(filepath))

        logger.info("[PPTX] 导出完成: %s, 共 %d 页", filepath, len(screenshots))
        return str(filepath)

    def _build_single_slide_html(self, slide_div, full_html: str) -> str:
        """构建单张幻灯片的独立 HTML"""
        soup = BeautifulSoup(full_html, "lxml")
        css = "\n".join(s.text for s in soup.find_all("style"))
        slide_html = str(slide_div)
        w, h = settings.PPT_PX_W, settings.PPT_PX_H

        return f"""<!DOCTYPE html>
<html lang="zh-CN">
<head>
<meta charset="UTF-8">
<style>
{css}
body {{
  margin: 0; padding: 0;
  display: flex; justify-content: center; align-items: center;
  width: {w}px; height: {h}px;
  overflow: hidden; background: {settings.PPT_BG_COLOR};
}}
.slide {{
  width: {w}px; height: {h}px;
  position: relative; overflow: hidden; flex-shrink: 0;
}}
</style>
</head>
<body>
{slide_html}
</body>
</html>"""

    def get_download_path(self, plan_id: str) -> str:
        filepath = self.upload_dir / f"ppt_{plan_id}.pptx"
        return str(filepath) if filepath.exists() else ""


ppt_export_service = PPTExportService()
